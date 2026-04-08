from docx import Document
from pptx import Presentation
from pathlib import Path

doc_path = Path(r"C:\Users\Qiang\.openclaw\media\outbound\42abaf42-c0fb-4d74-be24-460b98b5494c.docx")
ppt_path = Path(r"C:\Users\Qiang\.openclaw\media\outbound\65f80e1a-6c5e-4092-987a-204ff4f7df65.pptx")
out_doc = Path(r"C:\Users\Qiang\.openclaw\workspace-autosavants\tmp_assignment_template_extract.txt")
out_ppt = Path(r"C:\Users\Qiang\.openclaw\workspace-autosavants\tmp_business_model_ppt_extract.txt")

doc = Document(doc_path)
with out_doc.open('w', encoding='utf-8') as f:
    c=0
    for p in doc.paragraphs:
        t=(p.text or '').strip()
        if not t:
            continue
        c += 1
        f.write(f"{c:03d}. {t}\n")

prs = Presentation(ppt_path)
with out_ppt.open('w', encoding='utf-8') as f:
    for si, slide in enumerate(prs.slides, 1):
        f.write(f"\n--- Slide {si} ---\n")
        texts=[]
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                t=(shape.text or '').strip()
                if t:
                    texts.append(t)
        if not texts:
            f.write('(no text)\n')
        else:
            for t in texts:
                f.write(t.replace('\n',' / ') + '\n')
print('written')
