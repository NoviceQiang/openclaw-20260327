from docx import Document
from pptx import Presentation
from pathlib import Path

doc_path = Path(r"C:\Users\Qiang\.openclaw\media\outbound\42abaf42-c0fb-4d74-be24-460b98b5494c.docx")
ppt_path = Path(r"C:\Users\Qiang\.openclaw\media\outbound\65f80e1a-6c5e-4092-987a-204ff4f7df65.pptx")

print('=== DOCX TEMPLATE TEXT (first 220 non-empty paragraphs) ===')
doc = Document(doc_path)
c = 0
for i, p in enumerate(doc.paragraphs, 1):
    t = (p.text or '').strip()
    if not t:
        continue
    c += 1
    print(f'{c:03d}. {t}')
    if c >= 220:
        break

print('\n=== PPT TEXT (all slides) ===')
prs = Presentation(ppt_path)
for si, slide in enumerate(prs.slides, 1):
    print(f'\n--- Slide {si} ---')
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            t = (shape.text or '').strip()
            if t:
                texts.append(t)
    if not texts:
        print('(no text)')
    else:
        for t in texts:
            print(t.replace('\n', ' / '))
